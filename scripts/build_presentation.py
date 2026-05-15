"""Build the coursework defense presentation (.pptx).

Style notes:
- Neutral palette (white background, dark navy headings, dark gray body text,
  thin accent line).
- Sans-serif font (Calibri) per the guideline (Arial/Verdana/Tahoma class).
- Title >= 36pt, body >= 24pt.
- Bullets are concise and reflect the structure of the coursework.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# --- Palette -----------------------------------------------------------------
BG = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1F, 0x2D, 0x4E)   # deep navy
TEXT_COLOR = RGBColor(0x22, 0x2A, 0x35)    # near black
MUTED_COLOR = RGBColor(0x55, 0x5F, 0x70)
ACCENT_COLOR = RGBColor(0x2E, 0x5A, 0xA0)  # blue accent line
RULE_COLOR = RGBColor(0xD0, 0xD6, 0xE0)

FONT_NAME = "Calibri"

# --- Presentation setup ------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def _set_run(run, *, size=24, bold=False, color=TEXT_COLOR, italic=False,
             font=FONT_NAME):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_background(slide, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_header(slide, title_text, subtitle=None):
    """Header with title and a thin accent rule beneath."""
    title_box = slide.shapes.add_textbox(
        MARGIN_L, Inches(0.35), CONTENT_W, Inches(0.9)
    )
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    _set_run(r, size=32, bold=True, color=TITLE_COLOR)

    # Thin accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_L, Inches(1.25), Inches(1.4), Emu(28000)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR

    if subtitle:
        sb = slide.shapes.add_textbox(
            MARGIN_L, Inches(1.32), CONTENT_W, Inches(0.45)
        )
        stf = sb.text_frame
        stf.margin_left = 0
        stf.margin_top = 0.05
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        _set_run(sr, size=18, italic=True, color=MUTED_COLOR)


def add_footer(slide, page, total):
    fb = slide.shapes.add_textbox(
        MARGIN_L, SLIDE_H - Inches(0.45),
        CONTENT_W, Inches(0.35)
    )
    tf = fb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page} / {total}"
    _set_run(r, size=12, color=MUTED_COLOR)


def add_bullets(slide, items, *, left=None, top=None, width=None, height=None,
                size=22, line_spacing=1.25, color=TEXT_COLOR,
                first_bold=False):
    if left is None:
        left = MARGIN_L
    if top is None:
        top = Inches(1.75)
    if width is None:
        width = CONTENT_W
    if height is None:
        height = SLIDE_H - top - Inches(0.6)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, tail = item
        else:
            head, tail = None, item

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        # Bullet dash
        r0 = p.add_run()
        r0.text = "•  "
        _set_run(r0, size=size, color=ACCENT_COLOR, bold=True)

        if head:
            rh = p.add_run()
            rh.text = head
            _set_run(rh, size=size, color=color, bold=True)
            rs = p.add_run()
            rs.text = " — " + tail
            _set_run(rs, size=size, color=color)
        else:
            r1 = p.add_run()
            r1.text = tail
            _set_run(r1, size=size, color=color,
                     bold=first_bold and i == 0)

    return box


def add_paragraph_box(slide, text, *, left=None, top=None, width=None,
                      height=None, size=22, color=TEXT_COLOR,
                      align=PP_ALIGN.LEFT, bold=False, italic=False,
                      line_spacing=1.2):
    if left is None:
        left = MARGIN_L
    if top is None:
        top = Inches(1.75)
    if width is None:
        width = CONTENT_W
    if height is None:
        height = Inches(0.8)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    _set_run(r, size=size, color=color, bold=bold, italic=italic)
    return box


def add_card(slide, left, top, width, height, *, fill=RGBColor(0xF5, 0xF7, 0xFB),
             border=RULE_COLOR):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    card.text_frame.word_wrap = True
    return card


def style_table(table, *, header_fill=TITLE_COLOR, header_text=RGBColor(0xFF, 0xFF, 0xFF),
                row_fill_a=RGBColor(0xFF, 0xFF, 0xFF),
                row_fill_b=RGBColor(0xF3, 0xF5, 0xFA),
                size=16, header_size=16):
    rows = len(table.rows)
    cols = len(table.columns)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_NAME
                    if r == 0:
                        run.font.size = Pt(header_size)
                        run.font.bold = True
                        run.font.color.rgb = header_text
                    else:
                        run.font.size = Pt(size)
                        run.font.color.rgb = TEXT_COLOR
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_fill_a if (r % 2) else row_fill_b


# =============================================================================
# Slides
# =============================================================================

TOTAL = 12

# -------------------- Slide 1 : Title ----------------------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)

# left accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SLIDE_H)
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_COLOR

# Institution
inst = s.shapes.add_textbox(Inches(0.9), Inches(0.55),
                            SLIDE_W - Inches(1.5), Inches(0.9))
itf = inst.text_frame
itf.word_wrap = True
p = itf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "Министерство науки и высшего образования Российской Федерации"
_set_run(r, size=16, color=MUTED_COLOR)
p2 = itf.add_paragraph()
r2 = p2.add_run()
r2.text = "Образовательная организация высшего образования"
_set_run(r2, size=16, color=MUTED_COLOR)

# Theme block
theme_label = s.shapes.add_textbox(Inches(0.9), Inches(2.0),
                                   SLIDE_W - Inches(1.8), Inches(0.5))
tl = theme_label.text_frame
tl.margin_left = 0
pl = tl.paragraphs[0]
rl = pl.add_run()
rl.text = "КУРСОВАЯ РАБОТА"
_set_run(rl, size=20, bold=True, color=ACCENT_COLOR)

theme = s.shapes.add_textbox(Inches(0.9), Inches(2.5),
                             SLIDE_W - Inches(1.8), Inches(2.2))
ttf = theme.text_frame
ttf.word_wrap = True
ttf.margin_left = 0
p = ttf.paragraphs[0]
r = p.add_run()
r.text = ("Разработка программного модуля «Экономический календарь» "
         "для автоматизации сбора, хранения и визуализации "
         "макроэкономических событий")
_set_run(r, size=30, bold=True, color=TITLE_COLOR)

# Author / supervisor card
card_top = Inches(5.05)
card = add_card(s, Inches(0.9), card_top, SLIDE_W - Inches(1.8), Inches(1.6),
                fill=RGBColor(0xF5, 0xF7, 0xFB))
ctb = s.shapes.add_textbox(Inches(1.1), card_top + Inches(0.15),
                           SLIDE_W - Inches(2.2), Inches(1.35))
ctf = ctb.text_frame
ctf.word_wrap = True
ctf.margin_left = 0
for line in [
    ("Выполнил(а): ", "студент(ка) группы ____, ____________________"),
    ("Научный руководитель: ", "____________________"),
    ("Специальность / направление: ", "____________________"),
]:
    if line is None:
        continue
    p = ctf.add_paragraph() if ctf.paragraphs[0].runs else ctf.paragraphs[0]
    p.line_spacing = 1.25
    rl = p.add_run()
    rl.text = line[0]
    _set_run(rl, size=18, bold=True, color=TEXT_COLOR)
    rr = p.add_run()
    rr.text = line[1]
    _set_run(rr, size=18, color=TEXT_COLOR)

# City / year
cy = s.shapes.add_textbox(Inches(0.9), SLIDE_H - Inches(0.8),
                          SLIDE_W - Inches(1.8), Inches(0.5))
ctf = cy.text_frame
p = ctf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "2026"
_set_run(r, size=18, bold=True, color=MUTED_COLOR)


# -------------------- Slide 2 : Актуальность и проблема ----------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Актуальность и проблема исследования")

add_paragraph_box(
    s,
    "Макроэкономические события (решения по ставкам, ВВП, инфляция, занятость) "
    "напрямую влияют на финансовые рынки и требуют оперативного, "
    "структурированного доступа к данным.",
    top=Inches(1.55), height=Inches(1.1), size=20, color=TEXT_COLOR,
)

add_bullets(
    s,
    [
        ("Разрозненность источников", "данные распределены между календарями, "
         "новостями и аналитическими сервисами"),
        ("Ручная обработка", "ведение таблиц приводит к ошибкам и потере "
         "времени"),
        ("Низкая персонализация", "общий поток без учёта стран, типов и "
         "важности событий"),
        ("Перегруженность аналогов", "профессиональные платформы избыточны и "
         "дороги"),
    ],
    top=Inches(2.75),
    size=20,
)
add_footer(s, 2, TOTAL)


# -------------------- Slide 3 : Цель и задачи --------------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Цель и задачи работы")

# Goal card
goal = add_card(s, MARGIN_L, Inches(1.55), CONTENT_W, Inches(1.15),
                fill=RGBColor(0xEE, 0xF3, 0xFB))
gtb = s.shapes.add_textbox(MARGIN_L + Inches(0.25),
                           Inches(1.62),
                           CONTENT_W - Inches(0.5), Inches(1.0))
gtf = gtb.text_frame
gtf.word_wrap = True
gtf.margin_left = 0
p = gtf.paragraphs[0]
r = p.add_run()
r.text = "Цель: "
_set_run(r, size=22, bold=True, color=ACCENT_COLOR)
r2 = p.add_run()
r2.text = ("создание программного модуля для автоматизации работы с "
          "экономическим календарём, обеспечивающего удобный доступ к "
          "структурированной информации о событиях и их параметрах.")
_set_run(r2, size=22, color=TEXT_COLOR)

# Tasks
add_paragraph_box(
    s, "Задачи:",
    top=Inches(2.85), height=Inches(0.45), size=22, bold=True,
    color=TITLE_COLOR,
)
add_bullets(
    s,
    [
        "исследовать предметную область и проблемы существующих подходов",
        "выполнить сравнительный анализ аналогов и обосновать подход",
        "построить модели системы: Use Case, IDEF0, DFD, функциональную схему",
        "сформировать функциональные и нефункциональные требования",
        "обосновать архитектуру и технологический стек",
        "реализовать серверную и клиентскую части веб-приложения",
        "выполнить тестирование и оценить соответствие требованиям",
    ],
    top=Inches(3.35),
    size=18,
    line_spacing=1.15,
)
add_footer(s, 3, TOTAL)


# -------------------- Slide 4 : Анализ предметной области --------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 1. Анализ предметной области")

# Two columns: object / subject
col_w = (CONTENT_W - Inches(0.4)) / 2
col_top = Inches(1.55)
col_h = Inches(1.7)

c1 = add_card(s, MARGIN_L, col_top, col_w, col_h)
tb = s.shapes.add_textbox(MARGIN_L + Inches(0.2), col_top + Inches(0.15),
                          col_w - Inches(0.4), col_h - Inches(0.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Объект"
_set_run(r, size=18, bold=True, color=ACCENT_COLOR)
p2 = tf.add_paragraph(); p2.line_spacing = 1.2
r2 = p2.add_run()
r2.text = ("процесс управления и представления данных о макроэкономических "
          "событиях в информационных системах финансово-аналитического "
          "профиля")
_set_run(r2, size=17, color=TEXT_COLOR)

c2 = add_card(s, MARGIN_L + col_w + Inches(0.4), col_top, col_w, col_h)
tb2 = s.shapes.add_textbox(MARGIN_L + col_w + Inches(0.6),
                           col_top + Inches(0.15),
                           col_w - Inches(0.4), col_h - Inches(0.3))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "Предмет"
_set_run(r, size=18, bold=True, color=ACCENT_COLOR)
p2 = tf2.add_paragraph(); p2.line_spacing = 1.2
r2 = p2.add_run()
r2.text = ("разработка веб-приложения для автоматизированного сбора, "
          "хранения, фильтрации и визуализации данных экономического "
          "календаря с поддержкой мультиязычного интерфейса")
_set_run(r2, size=17, color=TEXT_COLOR)

add_paragraph_box(
    s, "Ключевые проблемы текущих подходов:",
    top=Inches(3.5), height=Inches(0.45), size=22, bold=True,
    color=TITLE_COLOR,
)
add_bullets(
    s,
    [
        "разрозненность источников и необходимость работы в нескольких сервисах",
        "временные затраты и ошибки при ручной обработке данных",
        "недостаточная персонализация и слабая система уведомлений",
        "перегруженность универсальных платформ непрофильным функционалом",
    ],
    top=Inches(4.0),
    size=18,
    line_spacing=1.15,
)
add_footer(s, 4, TOTAL)


# -------------------- Slide 5 : Сравнение аналогов ---------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 1. Сравнение существующих решений")

headers = ["Критерий", "TradingView", "Bloomberg Terminal", "Investing.com"]
rows_data = [
    ["Основной фокус", "Технический анализ, графики",
     "Полный спектр финансовых данных", "Макроэкономические события"],
    ["Экономический календарь", "Есть", "Есть", "Есть"],
    ["Уведомления", "Ограниченные", "Расширенные", "Push-уведомления"],
    ["Персонализация", "Ограниченная", "Высокая (проф.)", "Средняя"],
    ["Аналитика и прогнозы", "Технический анализ",
     "Глубокая фундаментальная", "Базовая"],
    ["Стоимость", "Бесплатно / подписка", "≈ $24 000 / год", "Бесплатно"],
    ["Перегруженность", "Высокая", "Очень высокая", "Средняя"],
]

t_left = MARGIN_L
t_top = Inches(1.55)
t_width = CONTENT_W
t_height = Inches(4.8)

table_shape = s.shapes.add_table(
    rows=1 + len(rows_data), cols=4,
    left=t_left, top=t_top, width=t_width, height=t_height
)
table = table_shape.table
# column widths
table.columns[0].width = Inches(2.7)
table.columns[1].width = Inches(3.2)
table.columns[2].width = Inches(3.2)
table.columns[3].width = Inches(3.0)

for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(rows_data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val

style_table(table, size=14, header_size=15)

add_paragraph_box(
    s,
    "Вывод: аналоги либо избыточно сложны и дороги, либо не специализированы "
    "под быстрый просмотр и фильтрацию событий — необходима собственная "
    "разработка.",
    top=Inches(6.55), height=Inches(0.6), size=15, italic=True,
    color=MUTED_COLOR,
)
add_footer(s, 5, TOTAL)


# -------------------- Slide 6 : Глава 2. Требования --------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 2. Требования к системе")

# Two columns
col_w = (CONTENT_W - Inches(0.4)) / 2
top = Inches(1.55)
height = Inches(5.2)

# Functional
add_card(s, MARGIN_L, top, col_w, height,
         fill=RGBColor(0xF5, 0xF7, 0xFB))
tb = s.shapes.add_textbox(MARGIN_L + Inches(0.25), top + Inches(0.15),
                          col_w - Inches(0.5), height - Inches(0.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Функциональные требования"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)

for item in [
    "просмотр календаря событий в виде списка",
    "автоматическое обновление данных",
    "поиск по ключевым словам",
    "детальная карточка события",
    "фильтрация и сортировка (страна, важность, период)",
    "раздел «Новости» из внешних источников",
]:
    pp = tf.add_paragraph()
    pp.line_spacing = 1.2
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = "•  "
    _set_run(r0, size=17, bold=True, color=ACCENT_COLOR)
    r1 = pp.add_run(); r1.text = item
    _set_run(r1, size=17, color=TEXT_COLOR)

# Non-functional
add_card(s, MARGIN_L + col_w + Inches(0.4), top, col_w, height,
         fill=RGBColor(0xF5, 0xF7, 0xFB))
tb2 = s.shapes.add_textbox(MARGIN_L + col_w + Inches(0.65),
                           top + Inches(0.15),
                           col_w - Inches(0.5), height - Inches(0.3))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "Нефункциональные требования"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)

for item in [
    "высокая отзывчивость интерфейса",
    "простая и предсказуемая навигация",
    "адаптивность под разные экраны",
    "поддержка мобильных устройств (Android 7.0+)",
    "устойчивость к сбоям сети за счёт кэширования",
    "мультиязычность интерфейса (RU/EN/ZH/ES)",
]:
    pp = tf2.add_paragraph()
    pp.line_spacing = 1.2
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = "•  "
    _set_run(r0, size=17, bold=True, color=ACCENT_COLOR)
    r1 = pp.add_run(); r1.text = item
    _set_run(r1, size=17, color=TEXT_COLOR)

add_footer(s, 6, TOTAL)


# -------------------- Slide 7 : Архитектура и стек ---------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 2. Архитектура и технологический стек")

# 3-tier architecture diagram (3 rounded boxes + arrows)
tiers = [
    ("Клиент", "React 18, Vite\nreact-router-dom\nЛокализация (i18n)\nТема (light/dark)"),
    ("Сервер", "FastAPI (Python)\nPydantic, REST API\nhttpx, Playwright\nUpsert-синхронизация"),
    ("Данные", "SQLAlchemy 2.0\nSQLite\nИндексы и external_id\nЛокальный кэш"),
]
box_top = Inches(1.7)
box_h = Inches(2.0)
gap = Inches(0.4)
box_w = (CONTENT_W - gap * 2) / 3

for i, (title, body) in enumerate(tiers):
    left = MARGIN_L + (box_w + gap) * i
    card = add_card(s, left, box_top, box_w, box_h,
                    fill=RGBColor(0xEE, 0xF3, 0xFB))
    tb = s.shapes.add_textbox(left + Inches(0.2), box_top + Inches(0.15),
                              box_w - Inches(0.4), box_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    _set_run(r, size=20, bold=True, color=ACCENT_COLOR)
    for line in body.split("\n"):
        pp = tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        pp.line_spacing = 1.15
        rr = pp.add_run(); rr.text = line
        _set_run(rr, size=15, color=TEXT_COLOR)

    if i < 2:
        arrow_left = left + box_w + Inches(0.05)
        arr = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            arrow_left, box_top + Inches(0.85),
            Inches(0.3), Inches(0.3)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = ACCENT_COLOR
        arr.line.fill.background()

# Stack justification
add_paragraph_box(
    s, "Обоснование выбора стека:",
    top=Inches(4.05), height=Inches(0.4),
    size=20, bold=True, color=TITLE_COLOR,
)
add_bullets(
    s,
    [
        ("FastAPI", "высокая производительность, async/await, валидация Pydantic, OpenAPI «из коробки»"),
        ("SQLAlchemy 2.0 + SQLite", "типизированная ORM, простота локального запуска, локальный кэш"),
        ("React 18 + Vite", "быстрая разработка SPA, удобная маршрутизация, мемоизация фильтров"),
        ("httpx + Playwright", "получение данных через JSON-API и fallback на HTML-таблицу"),
    ],
    top=Inches(4.5),
    size=16,
    line_spacing=1.15,
)
add_footer(s, 7, TOTAL)


# -------------------- Slide 8 : Реализация -----------------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 3. Реализация системы")

# Left: stages list
left_top = Inches(1.55)
left_h = Inches(5.2)
left_w = Inches(6.2)
add_card(s, MARGIN_L, left_top, left_w, left_h,
         fill=RGBColor(0xF5, 0xF7, 0xFB))
tb = s.shapes.add_textbox(MARGIN_L + Inches(0.25), left_top + Inches(0.15),
                          left_w - Inches(0.5), left_h - Inches(0.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Этапы разработки"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)

stages = [
    "проектирование модели данных Event (источник, важность, метрики)",
    "CRUD-слой и REST-эндпоинты на Pydantic-схемах",
    "синхронизация с внешними источниками и upsert по external_id",
    "фронтенд: таблица, календарь, фильтры, react-router-dom",
    "мультиязычность (RU/EN/ZH/ES) и темизация интерфейса",
    "интеграционная проверка сценариев API ↔ UI",
]
for i, item in enumerate(stages, start=1):
    pp = tf.add_paragraph()
    pp.line_spacing = 1.25
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = f"{i}.  "
    _set_run(r0, size=17, bold=True, color=ACCENT_COLOR)
    r1 = pp.add_run(); r1.text = item
    _set_run(r1, size=17, color=TEXT_COLOR)

# Right: key engineering features
right_left = MARGIN_L + left_w + Inches(0.3)
right_w = CONTENT_W - left_w - Inches(0.3)
add_card(s, right_left, left_top, right_w, left_h,
         fill=RGBColor(0xEE, 0xF3, 0xFB))
tb2 = s.shapes.add_textbox(right_left + Inches(0.25),
                           left_top + Inches(0.15),
                           right_w - Inches(0.5),
                           left_h - Inches(0.3))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "Ключевые решения"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)

features = [
    ("Upsert по external_id", "инкрементальная синхронизация без дублей"),
    ("Fallback JSON → HTML", "Playwright при недоступности API"),
    ("Локальный кэш SQLite", "работоспособность при сбоях источника"),
    ("useMemo на клиенте", "мгновенная фильтрация без запросов к серверу"),
    ("Авто-refresh при пустой выдаче", "повторная синхронизация при первом запуске"),
]
for head, tail in features:
    pp = tf2.add_paragraph()
    pp.line_spacing = 1.25
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = "•  "
    _set_run(r0, size=16, bold=True, color=ACCENT_COLOR)
    rh = pp.add_run(); rh.text = head
    _set_run(rh, size=16, bold=True, color=TEXT_COLOR)
    rt = pp.add_run(); rt.text = " — " + tail
    _set_run(rt, size=16, color=TEXT_COLOR)

add_footer(s, 8, TOTAL)


# -------------------- Slide 9 : Тестирование ---------------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 3. Тестирование и результаты проверки")

# Mini-summary
add_paragraph_box(
    s,
    "Тестирование выполнялось по 12 сценариям, охватывающим API, "
    "UI и устойчивость к сбоям внешних источников. Все сценарии пройдены.",
    top=Inches(1.5), height=Inches(0.7), size=18,
    color=MUTED_COLOR, italic=True,
)

# Table
headers = ["Код", "Сценарий", "Ожидаемый результат", "Итог"]
rows_data = [
    ["TC-01", "GET /health",
     "HTTP 200, {\"status\": \"ok\"}", "Пройден"],
    ["TC-02", "GET /events без фильтров",
     "Список из БД, сортировка по дате", "Пройден"],
    ["TC-03", "GET /events?country=US&importance=high",
     "Фильтрация по стране и важности", "Пройден"],
    ["TC-04", "POST /events/refresh",
     "Принудительный upsert и обновление", "Пройден"],
    ["TC-05/06", "GET /events/{id}/description?lang=ru|en",
     "Локализованное описание или fallback", "Пройден"],
    ["TC-07", "Модальное окно события",
     "Полные атрибуты события", "Пройден"],
    ["TC-08/09", "Переключение темы и языка",
     "Сохранение и применение настроек", "Пройден"],
    ["TC-10", "Календарное представление",
     "Группировка по датам и сортировка", "Пройден"],
    ["TC-11", "Отказ внешнего источника",
     "Работа за счёт кэша БД", "Пройден"],
    ["TC-12", "Повторная синхронизация",
     "Без дублей, обновление полей", "Пройден"],
]

t_top = Inches(2.3)
table_shape = s.shapes.add_table(
    rows=1 + len(rows_data), cols=4,
    left=MARGIN_L, top=t_top,
    width=CONTENT_W, height=Inches(4.7)
)
table = table_shape.table
table.columns[0].width = Inches(1.1)
table.columns[1].width = Inches(4.3)
table.columns[2].width = Inches(5.0)
table.columns[3].width = Inches(1.7)

for c, h in enumerate(headers):
    table.cell(0, c).text = h
for r, row in enumerate(rows_data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val

style_table(table, size=12, header_size=13)

# Color "Пройден" cells green
for r in range(1, len(rows_data) + 1):
    cell = table.cell(r, 3)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)
            run.font.bold = True

add_footer(s, 9, TOTAL)


# -------------------- Slide 10 : Демонстрация интерфейса ---------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Глава 3. Пользовательский интерфейс и сценарии")

# left: screens
add_paragraph_box(
    s, "Ключевые экраны:",
    top=Inches(1.55), height=Inches(0.45),
    size=20, bold=True, color=TITLE_COLOR,
)
add_bullets(
    s,
    [
        ("Главный экран «События»", "таблица с датой, временем, валютой, страной, важностью, прогнозом и фактом"),
        ("Карточка фильтров", "страна, период (сегодня/завтра/неделя/дата), важность, сброс"),
        ("Раздел «Новости»", "лента публикаций из внешних источников"),
        ("Модальное окно события", "описание и сравнение фактического / прогноза / предыдущего"),
        ("Окно чтения новости", "переход к источнику"),
        ("Переключатели", "язык (RU/EN/ZH/ES) и тема оформления"),
    ],
    top=Inches(2.05),
    width=Inches(6.6),
    size=15,
    line_spacing=1.15,
)

# right: scenarios
sc_left = MARGIN_L + Inches(6.8)
sc_w = CONTENT_W - Inches(6.8)
add_card(s, sc_left, Inches(1.55), sc_w, Inches(5.3),
         fill=RGBColor(0xF5, 0xF7, 0xFB))
tb = s.shapes.add_textbox(sc_left + Inches(0.25), Inches(1.7),
                          sc_w - Inches(0.5), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Сценарии использования"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)

scenarios = [
    ("1. Просмотр и фильтрация",
     "Открытие списка → автообновление → применение фильтров → мгновенная выборка."),
    ("2. Анализ события",
     "Выбор события → открытие карточки → сравнение факта/прогноза/предыдущего."),
    ("3. Новостной контекст",
     "Раздел «Новости» → окно чтения → переход к первоисточнику."),
]
for head, body in scenarios:
    pp = tf.add_paragraph()
    pp.line_spacing = 1.2
    pp.space_after = Pt(8)
    rh = pp.add_run(); rh.text = head
    _set_run(rh, size=16, bold=True, color=TEXT_COLOR)
    pp2 = tf.add_paragraph()
    pp2.line_spacing = 1.2
    pp2.space_after = Pt(6)
    rb = pp2.add_run(); rb.text = body
    _set_run(rb, size=15, color=MUTED_COLOR)

add_footer(s, 10, TOTAL)


# -------------------- Slide 11 : Выводы и развитие ---------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)
add_header(s, "Выводы и направления развития")

# Left: conclusions
col_w = (CONTENT_W - Inches(0.4)) / 2
top = Inches(1.55)
height = Inches(5.2)

add_card(s, MARGIN_L, top, col_w, height,
         fill=RGBColor(0xEE, 0xF3, 0xFB))
tb = s.shapes.add_textbox(MARGIN_L + Inches(0.25), top + Inches(0.15),
                          col_w - Inches(0.5), height - Inches(0.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Достигнутые результаты"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)
for item in [
    "проведён анализ предметной области и аналогов",
    "построены Use Case, IDEF0, DFD и функциональная схема",
    "сформированы функциональные и нефункциональные требования",
    "обоснован стек: FastAPI, SQLAlchemy 2.0, SQLite, React, Vite, Playwright",
    "реализованы клиент и сервер с авто-синхронизацией и кэшированием",
    "пройдены все 12 тестовых сценариев — цель работы достигнута",
]:
    pp = tf.add_paragraph()
    pp.line_spacing = 1.2
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = "•  "
    _set_run(r0, size=16, bold=True, color=ACCENT_COLOR)
    r1 = pp.add_run(); r1.text = item
    _set_run(r1, size=16, color=TEXT_COLOR)

# Right: future work
add_card(s, MARGIN_L + col_w + Inches(0.4), top, col_w, height,
         fill=RGBColor(0xF5, 0xF7, 0xFB))
tb2 = s.shapes.add_textbox(MARGIN_L + col_w + Inches(0.65),
                           top + Inches(0.15),
                           col_w - Inches(0.5), height - Inches(0.3))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
r = p.add_run(); r.text = "Направления развития"
_set_run(r, size=20, bold=True, color=ACCENT_COLOR)
for item in [
    "ролевая модель доступа и аутентификация",
    "уведомления: email, Telegram, push",
    "расширение перечня внешних источников",
    "переход на PostgreSQL и Alembic-миграции",
    "аналитические виджеты: тренды, корреляции",
    "автотесты (unit, integration, e2e) и CI-конвейер",
]:
    pp = tf2.add_paragraph()
    pp.line_spacing = 1.2
    pp.space_after = Pt(4)
    r0 = pp.add_run(); r0.text = "•  "
    _set_run(r0, size=16, bold=True, color=ACCENT_COLOR)
    r1 = pp.add_run(); r1.text = item
    _set_run(r1, size=16, color=TEXT_COLOR)

add_footer(s, 11, TOTAL)


# -------------------- Slide 12 : Спасибо за внимание -------------------------
s = prs.slides.add_slide(BLANK)
add_background(s)

# accent bar (right)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         SLIDE_W - Inches(0.35), 0,
                         Inches(0.35), SLIDE_H)
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_COLOR

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6),
                          SLIDE_W - Inches(1.8), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Спасибо за внимание!"
_set_run(r, size=54, bold=True, color=TITLE_COLOR)

tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.1),
                           SLIDE_W - Inches(1.8), Inches(0.8))
tf2 = tb2.text_frame
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Готов(а) ответить на ваши вопросы"
_set_run(r, size=26, color=MUTED_COLOR, italic=True)

# small note
tb3 = s.shapes.add_textbox(Inches(0.8), SLIDE_H - Inches(1.0),
                           SLIDE_W - Inches(1.8), Inches(0.5))
tf3 = tb3.text_frame
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Курсовая работа • Экономический календарь • 2026"
_set_run(r, size=14, color=MUTED_COLOR)


# Save
out = "/workspace/Презентация_курсовой_Экономический_календарь.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
